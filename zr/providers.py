from __future__ import annotations
import os, re, json, time
import pathlib
from typing import Any, Dict, Optional, List

DOI_RE = re.compile(r"\b10\.\d{4,9}/[-._;()/:A-Z0-9]+\b", re.I)

SCHEMA = {
  "type": "object",
  "properties": {
    "title":   {"type": ["string","null"]},
    "authors": {"type": ["array","null"], "items": {"type":"string"}},
    "year":    {"type": ["integer","null"]},
    "journal": {"type": ["string","null"]},
    "volume":  {"type": ["string","null"]},
    "issue":   {"type": ["string","null"]},
    "pages":   {"type": ["string","null"]},
    "doi":     {"type": ["string","null"]},
    "url":     {"type": ["string","null"]},
    "abstract":{"type": ["string","null"]},
    "record_type": {
      "type": ["string","null"],
      "enum": ["JOUR", "BOOK", "CHAP", "CONF", "RPRT", "THES", "WEB", "GEN", None],
      "description": "Document type: JOUR=Journal Article, BOOK=Book, CHAP=Book Chapter, CONF=Conference Paper, RPRT=Report, THES=Thesis, WEB=Web Page, GEN=Generic"
    },
    "evidence":{"type": "object"}
  },
  "required": ["title","authors","year","journal","doi","record_type","evidence"]
}

def _sanitize_jsonish(t: str) -> str:
    """
    Sanitize JSON-like text to make it parseable.
    Handles:
    - Markdown code blocks
    - Markdown-emphasized keys (*key*:)
    - Bare keys (key:)
    - Invalid control characters (ASCII 0-31 except \t, \n, \r)
    - Invalid Unicode escape sequences
    """
    if not t:
        return t

    # 去掉 ```json 代码块
    t = t.strip()
    t = re.sub(r"^```(?:json)?\s*", "", t, flags=re.I)
    t = re.sub(r"\s*```$", "", t)

    # 1) 把 *evidence*: 这种"markdown强调key"修成 "evidence":
    # 支持: *evidence*: / **evidence**:
    t = re.sub(r'(?m)^\s*\*+\s*([A-Za-z_][A-Za-z0-9_]*)\s*\*+\s*:', r'"\1":', t)

    # 2) 也顺手修一下常见的 "裸 key:" （可选，但很有用）
    # 只在行首匹配，避免伤到正文
    t = re.sub(r'(?m)^\s*([A-Za-z_][A-Za-z0-9_]*)\s*:', r'"\1":', t)

    # 3) 移除或转义无效的控制字符 (ASCII 0-31, 除了 \t=9, \n=10, \r=13)
    # JSON只允许这三个控制字符，其他的会导致解析失败
    def replace_control_char(match):
        char = match.group(0)
        code = ord(char)
        # 保留 tab, newline, carriage return
        if code in (9, 10, 13):
            return char
        # 其他控制字符：转义为 \uXXXX 或直接删除
        # 这里选择删除，因为它们通常是OCR错误或编码问题
        return ''

    t = re.sub(r'[\x00-\x1f]', replace_control_char, t)

    # 4) 修复 Python 风格的 \xXX 转义序列为 JSON 的 \uXXXX 格式
    # JSON 不支持 \xXX，需要转换为 \u00XX
    def fix_hex_escape(match):
        hex_part = match.group(1)  # 提取 XX 部分
        try:
            # 验证是否是有效的2位十六进制
            int(hex_part, 16)
            # 转换为 JSON Unicode 转义格式
            return f'\\u00{hex_part}'
        except ValueError:
            # 无效的十六进制，移除转义
            return hex_part

    # 匹配 \x 后跟 2 个字符
    t = re.sub(r'\\x([0-9a-fA-F]{2})', fix_hex_escape, t)

    # 5) 修复常见的无效 Unicode 转义序列
    # 例如: \u00e9 (有效) vs \u00E (无效，缺少一位) vs \u00rdi (无效，包含非十六进制字符)
    # 策略：移除无效的 \uXXXX 序列，保留有效的
    def fix_unicode_escape(match):
        escape_seq = match.group(0)
        # 检查是否是有效的4位十六进制
        hex_part = escape_seq[2:]  # 去掉 \u
        if len(hex_part) == 4:
            try:
                int(hex_part, 16)
                return escape_seq  # 有效，保留
            except ValueError:
                # 包含非十六进制字符，移除 \u 前缀，保留后面的字符
                return hex_part
        # 长度不足4位：移除整个转义序列
        return ''

    # 匹配 \u 后跟 0-4 个字符（可能是十六进制，也可能不是）
    # 使用非贪婪匹配，优先匹配4个字符
    t = re.sub(r'\\u([0-9a-fA-F]{0,4})', fix_unicode_escape, t)

    return t


def _extract_json_object(text: str) -> Optional[Dict[str, Any]]:
    """
    Extract and parse JSON object from text with multiple fallback strategies.

    Strategies (in order):
    1. Direct parsing (for clean JSON)
    2. Extract {...} and parse (for JSON with extra text)
    3. Sanitize and parse (for JSON with formatting issues)
    4. Sanitize extracted {...} and parse (for both issues)

    Returns:
        Dict if successful, None otherwise
    """
    if not text:
        return None
    t = text.strip()
    if not t:
        return None

    # Strategy 1: 先尝试直接解析（最常见的情况）
    try:
        obj = json.loads(t)
        return obj if isinstance(obj, dict) else None
    except json.JSONDecodeError:
        pass  # 继续尝试其他方法
    except Exception:
        pass

    # Strategy 2: 尝试：抽取最外层 {...}（处理有额外文本的情况）
    start = t.find("{")
    end = t.rfind("}")
    if start != -1 and end != -1 and end > start:
        cand = t[start:end+1]

        # 直接尝试解析提取的JSON
        try:
            obj = json.loads(cand)
            return obj if isinstance(obj, dict) else None
        except json.JSONDecodeError:
            pass  # 继续尝试修复
        except Exception:
            pass
    else:
        cand = t

    # Strategy 3: 对 JSON-ish 做一次"修复"再 parse（处理格式问题）
    cand2 = _sanitize_jsonish(cand)
    try:
        obj = json.loads(cand2)
        return obj if isinstance(obj, dict) else None
    except json.JSONDecodeError:
        pass  # 继续尝试最后一招
    except Exception:
        pass

    # Strategy 4: 如果上面都失败，尝试先sanitize整个文本，再提取JSON
    # 这处理了"有额外文本 + 格式问题"的组合情况
    if start != -1 and end != -1:
        t_sanitized = _sanitize_jsonish(t)
        start2 = t_sanitized.find("{")
        end2 = t_sanitized.rfind("}")
        if start2 != -1 and end2 != -1 and end2 > start2:
            cand3 = t_sanitized[start2:end2+1]
            try:
                obj = json.loads(cand3)
                return obj if isinstance(obj, dict) else None
            except Exception:
                pass

    return None


class BaseProvider:
    def __init__(self, model: str, base_url: str, timeout: int, num_ctx: int, api_key_env: str, trust_env: bool, log):
        self.model = model
        self.base_url = base_url
        self.timeout = timeout
        self.num_ctx = num_ctx
        self.api_key_env = api_key_env
        self.trust_env = trust_env
        self.log = log

    # ---------- PDF helpers ----------
    def pdf_text_snippet(self, pdf_path: pathlib.Path, max_pages: int, scan_doi_pages: int) -> Dict[str, Any]:
        import fitz  # PyMuPDF

        # Check if file is empty or corrupted before opening
        try:
            file_size = pdf_path.stat().st_size
            if file_size == 0:
                self.log.warning(f"PDF file is empty (0 bytes): {pdf_path}")
                return {"text": f"[EMPTY FILE: {pdf_path.name}]", "doi_hit": False}
        except Exception as e:
            self.log.warning(f"Cannot stat PDF file {pdf_path}: {e}")
            return {"text": f"[INACCESSIBLE FILE: {pdf_path.name}]", "doi_hit": False}

        try:
            doc = fitz.open(pdf_path)
        except fitz.EmptyFileError:
            self.log.warning(f"PDF file is empty or corrupted: {pdf_path}")
            return {"text": f"[EMPTY/CORRUPTED FILE: {pdf_path.name}]", "doi_hit": False}
        except Exception as e:
            self.log.warning(f"Cannot open PDF file {pdf_path}: {e}")
            return {"text": f"[UNREADABLE FILE: {pdf_path.name}]", "doi_hit": False}

        texts = []
        for i in range(min(max_pages, doc.page_count)):
            texts.append(doc.load_page(i).get_text("text"))

        doi_page_text = ""
        doi_hit = False
        for i in range(min(scan_doi_pages, doc.page_count)):
            t = doc.load_page(i).get_text("text")
            if DOI_RE.search(t):
                doi_page_text = t
                doi_hit = True
                break

        full = "\n\n".join([*texts, "\n\n[DOI_PAGE]\n" + (doi_page_text or "")])
        return {"text": full, "doi_hit": doi_hit}

    def quick_doi_hit(self, pdf_path: pathlib.Path, scan_pages: int) -> bool:
        try:
            import fitz
            # Check file size first to avoid EmptyFileError
            if pdf_path.stat().st_size == 0:
                return False
            doc = fitz.open(pdf_path)
            for i in range(min(scan_pages, doc.page_count)):
                t = doc.load_page(i).get_text("text")
                if DOI_RE.search(t):
                    return True
        except Exception:
            return False
        return False

    # ---------- Generic snippet ----------
    def generic_text_snippet(self, path: pathlib.Path) -> str:
        suf = path.suffix.lower()
        try:
            if suf == ".txt" or suf == ".md":
                return path.read_text(encoding="utf-8", errors="ignore")[:20000]

            # Handle old formats (.doc, .ppt) by converting to new formats first
            if suf in [".doc", ".ppt"]:
                converted_path = self._convert_old_office_format(path)
                if converted_path and converted_path.exists():
                    return self.generic_text_snippet(converted_path)
                else:
                    self.log.warning(f"Failed to convert old format {suf}: {path}")
                    return ""

            if suf == ".docx":
                from docx import Document
                doc = Document(str(path))
                parts = []

                # Extract document title and core metadata
                # Title is often in the first few paragraphs or in document properties
                if doc.core_properties.title:
                    parts.append(f"[DOCUMENT TITLE]: {doc.core_properties.title}")
                if doc.core_properties.author:
                    parts.append(f"[AUTHOR]: {doc.core_properties.author}")
                if doc.core_properties.created:
                    parts.append(f"[DATE]: {doc.core_properties.created.year}")

                parts.append("\n[DOCUMENT CONTENT]:")

                # Extract paragraphs with style information
                # Heading styles often contain important metadata
                for i, p in enumerate(doc.paragraphs[:200]):
                    text = p.text.strip()
                    if not text:
                        continue

                    # Mark headings/titles for better LLM understanding
                    if p.style.name.startswith('Heading') or p.style.name == 'Title':
                        parts.append(f"\n[HEADING]: {text}")
                    else:
                        parts.append(text)

                    # First 10 paragraphs are most important for metadata
                    if i < 10:
                        parts.append("")  # Add spacing for clarity

                return "\n".join(parts)[:20000]

            if suf == ".pptx":
                from pptx import Presentation
                prs = Presentation(str(path))
                parts = []

                # Convert slides to list to avoid iteration issues
                slide_list = list(prs.slides)

                # First slide usually contains title, author, date
                if len(slide_list) > 0:
                    first_slide = slide_list[0]
                    parts.append("[FIRST SLIDE - Title/Metadata]:")

                    for shape in first_slide.shapes:
                        if hasattr(shape, "text") and shape.text.strip():
                            text = shape.text.strip()
                            # Title placeholder is usually the main title
                            if hasattr(shape, "is_placeholder") and shape.is_placeholder:
                                parts.append(f"[TITLE]: {text}")
                            else:
                                parts.append(text)
                    parts.append("")

                # Extract subsequent slides with their titles
                parts.append("[SLIDE CONTENTS]:")
                for i in range(1, min(20, len(slide_list))):  # Process up to 20 slides
                    slide = slide_list[i]
                    slide_parts = []

                    for shape in slide.shapes:
                        if hasattr(shape, "text") and shape.text.strip():
                            text = shape.text.strip()
                            # First text box is often the slide title
                            if not slide_parts:
                                slide_parts.append(f"\n[SLIDE {i+1} TITLE]: {text}")
                            else:
                                slide_parts.append(text)

                    if slide_parts:
                        parts.extend(slide_parts)

                return "\n".join(parts)[:20000]
        except Exception:
            return ""
        return ""

    def _convert_old_office_format(self, path: pathlib.Path) -> Optional[pathlib.Path]:
        """
        Convert old Office formats (.doc, .ppt) to new formats (.docx, .pptx) using LibreOffice.
        Returns path to converted file, or None if conversion failed.
        """
        import subprocess
        import tempfile

        suf = path.suffix.lower()
        if suf not in [".doc", ".ppt"]:
            return None

        # Determine target format
        target_format = "docx" if suf == ".doc" else "pptx"

        # Create temp directory for converted file
        temp_dir = pathlib.Path(tempfile.gettempdir()) / "zotero_restore_converted"
        temp_dir.mkdir(parents=True, exist_ok=True)

        # Output path
        output_path = temp_dir / f"{path.stem}.{target_format}"

        # If already converted, return cached version
        if output_path.exists():
            return output_path

        # Try unoconv first (simpler)
        try:
            result = subprocess.run(
                ["unoconv", "-f", target_format, "-o", str(output_path), str(path)],
                capture_output=True,
                text=True,
                timeout=60
            )
            if result.returncode == 0 and output_path.exists():
                self.log.info(f"Converted {path.name} to {target_format} using unoconv")
                return output_path
        except FileNotFoundError:
            pass  # unoconv not installed, try libreoffice
        except Exception as e:
            self.log.warning(f"unoconv conversion failed: {e}")

        # Try LibreOffice directly
        try:
            result = subprocess.run(
                ["libreoffice", "--headless", "--convert-to", target_format,
                 "--outdir", str(temp_dir), str(path)],
                capture_output=True,
                text=True,
                timeout=60
            )
            if result.returncode == 0 and output_path.exists():
                self.log.info(f"Converted {path.name} to {target_format} using LibreOffice")
                return output_path
        except FileNotFoundError:
            self.log.warning("LibreOffice not found. Install LibreOffice or unoconv to support .doc/.ppt files.")
        except Exception as e:
            self.log.warning(f"LibreOffice conversion failed: {e}")

        return None

    # ---------- LLM calls (to be implemented) ----------
    def extract_biblio(self, text: str) -> Dict[str, Any]:
        raise NotImplementedError

    def extract_generic(self, text: str) -> Dict[str, Any]:
        # For non-paper docs: allow journal/doi null
        return self.extract_biblio(text)

    def choose_main_pdf(self, pdfs: List[pathlib.Path], max_pages: int, scan_pages: int) -> Optional[pathlib.Path]:
        # default: fallback only
        return None


class OllamaProvider(BaseProvider):
    def extract_biblio(self, text: str) -> Dict[str, Any]:
        import requests
        prompt = (
            "你是文献信息抽取器，只能基于我提供的文本抽取元数据，不要猜测。\n"
            "若字段不存在请填 null。\n"
            "authors 输出为字符串数组。\n"
            "doi 仅在文本中找到时填写。\n"
            "abstract 若文本中出现摘要则提取，否则为 null。\n"
            "record_type 根据内容判断文档类型：\n"
            "  * JOUR (期刊文章): 有期刊名、卷号、期号、DOI、摘要\n"
            "  * BOOK (书籍): 书名、出版社、ISBN、章节\n"
            "  * CHAP (书籍章节): 有书名和章节标题\n"
            "  * CONF (会议论文): 会议论文集\n"
            "  * RPRT (报告): 技术报告、白皮书\n"
            "  * THES (学位论文): 硕士或博士论文\n"
            "  * WEB (网页): 网页内容（有URL，无期刊/出版社）\n"
            "  * GEN (通用): 演示文稿、笔记等\n"
            "同时在 evidence 里给每个字段提供对应的原文短片段。\n"
            f"JSON schema: {json.dumps(SCHEMA, ensure_ascii=False)}\n"
            "文本如下：\n" + (text or "")
        )
        payload = {
            "model": self.model,
            "messages": [{"role": "user", "content": prompt}],
            "format": SCHEMA,
            "stream": False,
            "options": {"num_ctx": self.num_ctx}
        }
        url = self.base_url.rstrip("/") + "/api/chat"
        r = requests.post(url, json=payload, timeout=self.timeout)
        r.raise_for_status()
        content = r.json()["message"]["content"]
        if isinstance(content, dict):
            return content
        obj = _extract_json_object(content)
        if not obj:
            raise ValueError("Failed to parse JSON from ollama response.")
        return obj

    def choose_main_pdf(self, pdfs: List[pathlib.Path], max_pages: int, scan_pages: int) -> Optional[pathlib.Path]:
        # Minimal LLM chooser using ollama structured output: {"index": int}
        import requests
        blocks = []
        for i, p in enumerate(pdfs, 1):
            info = self.pdf_text_snippet(p, max_pages=max_pages, scan_doi_pages=scan_pages)
            snip = (info["text"] or "").replace("\n", " ").strip()[:2500]
            blocks.append(f"[{i}] {p.name}\n{snip}\n")
        schema = {
            "type": "object",
            "properties": {"index": {"type": "integer"}},
            "required": ["index"]
        }
        prompt = (
            "同一文件夹内有多个包含DOI的PDF。请选择最像“论文正文/主文(full text)”的那个，"
            "而不是supplement/appendix/supporting information。\n"
            "只输出 JSON: {\"index\": 1} 这种。\n\n" + "\n".join(blocks)
        )
        payload = {
            "model": self.model,
            "messages": [{"role": "user", "content": prompt}],
            "format": schema,
            "stream": False,
            "options": {"num_ctx": self.num_ctx}
        }
        url = self.base_url.rstrip("/") + "/api/chat"
        r = requests.post(url, json=payload, timeout=self.timeout)
        if r.status_code != 200:
            return None
        content = r.json()["message"]["content"]
        if isinstance(content, str):
            obj = _extract_json_object(content)
        else:
            obj = content
        try:
            idx = int(obj["index"])
            if 1 <= idx <= len(pdfs):
                return pdfs[idx-1]
        except Exception:
            return None
        return None


class OpenAISDKProvider(BaseProvider):
    def extract_biblio(self, text: str) -> Dict[str, Any]:
        from openai import OpenAI
        import httpx

        api_key = os.getenv(self.api_key_env) or os.getenv("OPENAI_API_KEY")
        if not api_key:
            raise RuntimeError(f"Missing API key env: {self.api_key_env} (or OPENAI_API_KEY)")

        client = OpenAI(
            api_key=api_key,
            base_url=self.base_url,
            timeout=self.timeout,
            http_client=httpx.Client(trust_env=self.trust_env),
        )

        schema_str = json.dumps(SCHEMA, ensure_ascii=False)

        system = (
            "You are a bibliographic metadata extractor. "
            "CRITICAL: Return ONLY a valid JSON object. "
            "No markdown code blocks, no explanations, no extra text. "
            "The JSON must be parseable by json.loads() without any preprocessing."
        )

        prompt = (
            "Extract metadata ONLY from the provided text. Do not guess or invent information.\n"
            "Rules:\n"
            "- If a field is missing, use null\n"
            "- authors must be a JSON array of strings (e.g., [\"Author1\", \"Author2\"])\n"
            "- abstract: extract if present in text, otherwise null\n"
            "- doi: only fill if found in text\n"
            "- record_type: Identify document type based on content:\n"
            "  * JOUR (Journal Article): Has journal name, volume, issue, DOI, abstract\n"
            "  * BOOK: Book title, publisher, ISBN, chapters\n"
            "  * CHAP: Book chapter with book title and chapter title\n"
            "  * CONF: Conference paper/proceedings\n"
            "  * RPRT: Technical report, white paper\n"
            "  * THES: Thesis or dissertation\n"
            "  * WEB: Web page content (has URL, no journal/publisher)\n"
            "  * GEN: Generic document (presentations, notes, etc.)\n"
            "- evidence: provide a dict with field names as keys\n"
            "- Return ONLY the JSON object, no markdown, no extra text\n\n"
            f"JSON schema:\n{schema_str}\n\n"
            f"TEXT:\n{text}"
        )

        def dump_raw(tag: str, s: str) -> pathlib.Path:
            log_dir = pathlib.Path(os.getenv("LOG_DIR", "./logs")).resolve()
            log_dir.mkdir(parents=True, exist_ok=True)
            p = log_dir / f"{tag}_{int(time.time())}.txt"
            p.write_text(s or "", encoding="utf-8", errors="ignore")
            return p

        # ---------- 1st call ----------
        resp = client.chat.completions.create(
            model=self.model,
            messages=[
                {"role": "system", "content": system},
                {"role": "user", "content": prompt},
            ],
            temperature=0,
        )
        content = (resp.choices[0].message.content or "").strip()
        obj = _extract_json_object(content)
        if obj:
            return obj

        raw1 = dump_raw("openai_raw_1", content)
        self.log.warning(f"OpenAI output not JSON. raw saved: {raw1}")

        # ---------- 2nd call (stricter) ----------
        retry = (
            "IMPORTANT: Return ONLY one JSON object that matches the schema. "
            "No markdown. No extra text.\n"
            f"JSON schema:\n{schema_str}\n\n"
            f"TEXT:\n{text}"
        )
        resp2 = client.chat.completions.create(
            model=self.model,
            messages=[{"role": "system", "content": system}, {"role": "user", "content": retry}],
            temperature=0,
        )
        content2 = (resp2.choices[0].message.content or "").strip()
        obj2 = _extract_json_object(content2)
        if obj2:
            return obj2

        raw2 = dump_raw("openai_raw_2", content2)
        self.log.warning(f"OpenAI output not JSON. raw saved: {raw2}")

        # ---------- 3rd call: JSON repair ----------
        repair_system = "You are a JSON repair tool. Output ONLY valid JSON. No markdown. No extra text."
        repair_user = (
            "Fix the following output into a SINGLE valid JSON object that conforms to the given schema.\n\n"
            f"JSON schema:\n{schema_str}\n\n"
            "RAW OUTPUT:\n"
            f"{content2}"
        )
        resp3 = client.chat.completions.create(
            model=self.model,
            messages=[
                {"role": "system", "content": repair_system},
                {"role": "user", "content": repair_user},
            ],
            temperature=0,
        )
        content3 = (resp3.choices[0].message.content or "").strip()
        obj3 = _extract_json_object(content3)
        if obj3:
            return obj3

        raw3 = dump_raw("openai_raw_repair", content3)
        raise ValueError(f"Failed to parse JSON from OpenAI SDK response even after repair. raw saved: {raw3}")


    def choose_main_pdf(self, pdfs: List[pathlib.Path], max_pages: int, scan_pages: int) -> Optional[pathlib.Path]:
        from openai import OpenAI
        import httpx
        api_key = os.getenv(self.api_key_env) or os.getenv("OPENAI_API_KEY")
        if not api_key:
            return None

        client = OpenAI(
            api_key=api_key,
            base_url=self.base_url,
            timeout=self.timeout,
            http_client=httpx.Client(trust_env=self.trust_env),
        )

        blocks = []
        for i, p in enumerate(pdfs, 1):
            info = self.pdf_text_snippet(p, max_pages=max_pages, scan_doi_pages=scan_pages)
            snip = (info["text"] or "").replace("\n", " ").strip()[:2500]
            blocks.append(f"[{i}] file={p.name}\ntext={snip}\n")

        prompt = (
            "In one folder there are multiple PDFs with DOI. Choose which is the MAIN/FULL-TEXT paper, "
            "NOT supplementary/appendix/supporting info.\n"
            "Return ONLY the number (e.g., 1).\n\n" + "\n".join(blocks)
        )

        resp = client.chat.completions.create(
            model=self.model,
            messages=[{"role": "user", "content": prompt}],
            temperature=0,
        )
        out = (resp.choices[0].message.content or "").strip()
        m = re.search(r"\d+", out)
        if not m:
            return None
        idx = int(m.group(0))
        if 1 <= idx <= len(pdfs):
            return pdfs[idx-1]
        return None


def build_provider(provider_name: str, model: str, base_url: str, timeout: int, num_ctx: int,
                   api_key_env: str, trust_env: bool, log):
    if provider_name == "ollama_native":
        return OllamaProvider(model, base_url, timeout, num_ctx, api_key_env, trust_env, log)
    if provider_name == "openai_sdk":
        return OpenAISDKProvider(model, base_url, timeout, num_ctx, api_key_env, trust_env, log)
    raise ValueError(f"Unknown provider: {provider_name}")

